"""PPT Generator機能のビジネスロジックサービス。

このモジュールは、PPTファイルの生成・ダウンロード・画像エクスポート・
スライド選択エクスポートなどのビジネスロジックを提供します。

主な機能:
    - PPTファイルのダウンロード
    - 選択されたスライドのみをエクスポート
    - スライドを画像として取得
    - 質問データ（Excel→CSV）のダウンロード
    - PPTファイルのアップロード

使用例:
    >>> from app.services.ppt_generator import PPTGeneratorService
    >>>
    >>> ppt_service = PPTGeneratorService()
    >>> content = await ppt_service.download_ppt(
    ...     package="market-analysis",
    ...     phase="phase1",
    ...     template="template_a"
    ... )
"""

import os
from io import BytesIO

import pandas as pd
from pptx import Presentation

from app.api.decorators import async_timeout, measure_performance
from app.core.exceptions import NotFoundError, ValidationError
from app.core.logging import get_logger
from app.services.storage import StorageService, get_storage_service

logger = get_logger(__name__)

# PPT Generator用のストレージ設定
PPT_STORAGE_CONTAINER = "ppt-generator"


class PPTGeneratorService:
    """PPT Generator機能のビジネスロジックを提供するサービスクラス。

    このサービスは、PPTファイルの操作に関するすべての操作を提供します。
    すべての操作は非同期で実行され、適切なロギングとエラーハンドリングを含みます。

    Attributes:
        storage: StorageServiceインスタンス（ファイルストレージ用）
        container: コンテナ名（デフォルト: "ppt-generator"）

    Example:
        >>> ppt_service = PPTGeneratorService()
        >>> content = await ppt_service.download_ppt("pkg", "phase1", "temp_a")
    """

    def __init__(self, storage: StorageService | None = None):
        """PPT Generatorサービスを初期化します。

        Args:
            storage (StorageService | None): ストレージサービスインスタンス
                - Noneの場合、get_storage_service()で自動取得

        Note:
            - テスト時はモックストレージを注入可能
        """
        self.storage = storage or get_storage_service()
        self.container = PPT_STORAGE_CONTAINER
        logger.info(
            "PPT Generatorサービスを初期化しました",
            container=self.container,
        )

    def _generate_file_path(self, package: str, phase: str, template: str, filename: str = "") -> str:
        """ストレージパスを生成します。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            filename: ファイル名（オプション）

        Returns:
            str: ストレージパス（例: "package/phase/template/file.pptx"）

        Example:
            >>> path = self._generate_file_path("pkg", "phase1", "temp_a", "slide.pptx")
            >>> print(path)
            pkg/phase1/temp_a/slide.pptx
        """
        if filename:
            return f"{package}/{phase}/{template}/{filename}"
        return f"{package}/{phase}/{template}"

    @measure_performance
    @async_timeout(seconds=180)  # 3分タイムアウト
    async def find_file_in_directory(self, package: str, phase: str, template: str, extension: str = "pptx") -> str:
        """ディレクトリ内から指定された拡張子のファイルを検索します。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            extension: ファイル拡張子（デフォルト: "pptx"）

        Returns:
            str: ファイルのストレージパス

        Raises:
            NotFoundError: ファイルが見つからない場合

        Example:
            >>> path = await service.find_file_in_directory("pkg", "phase1", "temp_a")
            >>> print(path)
            pkg/phase1/temp_a/presentation.pptx

        Note:
            - 最初に見つかったファイルを返します
            - 複数ファイルがある場合は最初のものが選択されます
        """
        logger.debug(
            "ファイル検索中",
            package=package,
            phase=phase,
            template=template,
            extension=extension,
            action="find_file_in_directory",
        )

        try:
            target_dir = self._generate_file_path(package, phase, template)

            # ディレクトリ内のファイル一覧を取得
            file_list = await self.storage.list_blobs(self.container, prefix=target_dir)

            # 指定された拡張子でフィルタ
            matching_files = [f for f in file_list if f.endswith(f".{extension}")]

            if not matching_files:
                logger.warning(
                    "ファイルが見つかりません",
                    package=package,
                    phase=phase,
                    template=template,
                    extension=extension,
                )
                raise NotFoundError(
                    f"指定されたディレクトリに{extension}ファイルが見つかりません",
                    details={
                        "package": package,
                        "phase": phase,
                        "template": template,
                        "extension": extension,
                    },
                )

            # 最初に見つかったファイルを返す
            file_path = matching_files[0]

            logger.debug(
                "ファイルを見つけました",
                file_path=file_path,
            )

            return file_path

        except (NotFoundError, ValidationError):
            raise
        except Exception as e:
            logger.error(
                "ファイル検索中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "ファイル検索に失敗しました",
                details={"error": str(e)},
            ) from e

    @measure_performance
    @async_timeout(seconds=300)  # 5分タイムアウト
    async def download_ppt(self, package: str, phase: str, template: str) -> bytes:
        """PPTファイルをダウンロードします。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名

        Returns:
            bytes: PPTファイルのバイナリデータ

        Raises:
            NotFoundError: ファイルが存在しない場合
            ValidationError: ダウンロードに失敗した場合

        Example:
            >>> content = await service.download_ppt("pkg", "phase1", "temp_a")
            >>> print(f"Downloaded {len(content)} bytes")
            Downloaded 1024000 bytes

        Note:
            - タイムアウトは5分です
            - 大きなファイルはタイムアウトする可能性があります
        """
        logger.info(
            "PPTファイルダウンロード中",
            package=package,
            phase=phase,
            template=template,
            action="download_ppt",
        )

        try:
            # ファイルパスを検索
            file_path = await self.find_file_in_directory(package, phase, template, extension="pptx")

            # ファイルをダウンロード
            content = await self.storage.download(self.container, file_path)

            logger.info(
                "PPTファイルをダウンロードしました",
                file_path=file_path,
                size=len(content),
            )

            return content

        except (NotFoundError, ValidationError):
            raise
        except Exception as e:
            logger.error(
                "PPTファイルダウンロード中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "PPTファイルのダウンロードに失敗しました",
                details={"error": str(e)},
            ) from e

    @measure_performance
    @async_timeout(seconds=300)  # 5分タイムアウト
    async def export_selected_slides(self, package: str, phase: str, template: str, slide_numbers: str) -> bytes:
        """選択されたスライドのみを含む新しいPPTファイルを生成します。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            slide_numbers: スライド番号のカンマ区切り文字列（例: "1,3,5,7"）

        Returns:
            bytes: 選択されたスライドのみを含むPPTファイルのバイナリデータ

        Raises:
            NotFoundError: ファイルが存在しない場合
            ValidationError: エクスポートに失敗した場合

        Example:
            >>> content = await service.export_selected_slides(
            ...     "pkg", "phase1", "temp_a", "1,3,5,7"
            ... )
            >>> print(f"Exported {len(content)} bytes")
            Exported 512000 bytes

        Note:
            - スライド番号は1から開始します
            - 存在しないスライド番号は無視されます
            - タイムアウトは5分です
        """
        logger.info(
            "選択されたスライドをエクスポート中",
            package=package,
            phase=phase,
            template=template,
            slide_numbers=slide_numbers,
            action="export_selected_slides",
        )

        try:
            # 元のPPTファイルをダウンロード
            file_path = await self.find_file_in_directory(package, phase, template, extension="pptx")
            pptx_bytes = await self.storage.download(self.container, file_path)

            # PPTXファイルをロード
            prs = Presentation(BytesIO(pptx_bytes))

            # 選択されたスライド番号をパース（1-indexed → 0-indexed）
            selected_indices = {int(num.strip()) - 1 for num in slide_numbers.split(",")}

            logger.debug(
                "スライドをフィルタリング中",
                total_slides=len(prs.slides),
                selected_slides=len(selected_indices),
            )

            # 選択されていないスライドを削除（逆順で削除）
            total_slides = len(prs.slides)
            for i in reversed(range(total_slides)):
                if i not in selected_indices:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]

            # 新しいPPTXファイルを生成
            output_stream = BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)
            content = output_stream.read()

            logger.info(
                "選択されたスライドをエクスポートしました",
                original_slides=total_slides,
                exported_slides=len(selected_indices),
                size=len(content),
            )

            return content

        except (NotFoundError, ValidationError):
            raise
        except Exception as e:
            logger.error(
                "スライドエクスポート中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                slide_numbers=slide_numbers,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "スライドのエクスポートに失敗しました",
                details={"error": str(e)},
            ) from e

    @measure_performance
    @async_timeout(seconds=180)  # 3分タイムアウト
    async def get_slide_image(self, package: str, phase: str, template: str, slide_number: int) -> bytes:
        """スライドを画像として取得します。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            slide_number: スライド番号（1から開始）

        Returns:
            bytes: PNG画像のバイナリデータ

        Raises:
            NotFoundError: 画像が存在しない場合
            ValidationError: 取得に失敗した場合

        Example:
            >>> content = await service.get_slide_image(
            ...     "pkg", "phase1", "temp_a", 5
            ... )
            >>> print(f"Downloaded {len(content)} bytes")
            Downloaded 102400 bytes

        Note:
            - 画像は事前に生成されている必要があります
            - 画像パス: {package}/{phase}/{template}/img/XXX.png
            - タイムアウトは3分です
        """
        logger.info(
            "スライド画像取得中",
            package=package,
            phase=phase,
            template=template,
            slide_number=slide_number,
            action="get_slide_image",
        )

        try:
            # 画像パスを生成（3桁ゼロパディング）
            image_path = self._generate_file_path(package, phase, template, f"img/{slide_number:03d}.png")

            # 画像をダウンロード
            content = await self.storage.download(self.container, image_path)

            logger.info(
                "スライド画像を取得しました",
                image_path=image_path,
                size=len(content),
            )

            return content

        except (NotFoundError, ValidationError):
            raise
        except Exception as e:
            logger.error(
                "スライド画像取得中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                slide_number=slide_number,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "スライド画像の取得に失敗しました",
                details={"error": str(e)},
            ) from e

    @measure_performance
    @async_timeout(seconds=300)  # 5分タイムアウト
    async def download_question(self, package: str, phase: str, template: str, question_type: str) -> tuple[bytes, str]:
        """質問データをExcelからCSV形式でダウンロードします。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            question_type: 質問タイプ（Excelシート名）

        Returns:
            tuple[bytes, str]: (CSVデータ, ファイル名) のタプル

        Raises:
            NotFoundError: ファイルが存在しない場合
            ValidationError: ダウンロードに失敗した場合

        Example:
            >>> content, filename = await service.download_question(
            ...     "pkg", "phase1", "temp_a", "customer_survey"
            ... )
            >>> print(filename)
            questions_customer_survey.csv

        Note:
            - ExcelファイルからCSVに変換されます
            - UTF-8 BOMエンコーディングで出力されます
            - タイムアウトは5分です
        """
        logger.info(
            "質問データダウンロード中",
            package=package,
            phase=phase,
            template=template,
            question_type=question_type,
            action="download_question",
        )

        try:
            # Excelファイルを検索
            file_path = await self.find_file_in_directory(package, phase, template, extension="xlsx")

            # Excelファイルをダウンロード
            excel_bytes = await self.storage.download(self.container, file_path)

            # 指定されたシートを読み込み
            df = pd.read_excel(BytesIO(excel_bytes), sheet_name=question_type)

            logger.debug(
                "Excelシートを読み込みました",
                sheet_name=question_type,
                rows=len(df),
                columns=len(df.columns),
            )

            # CSVに変換
            output_stream = BytesIO()
            df.to_csv(output_stream, index=False, encoding="utf-8-sig")
            output_stream.seek(0)
            content = output_stream.read()

            # ファイル名を生成
            file_name = os.path.basename(file_path)
            name, _ = os.path.splitext(file_name)
            csv_filename = f"{name}_{question_type}.csv"

            logger.info(
                "質問データをダウンロードしました",
                csv_filename=csv_filename,
                size=len(content),
            )

            return content, csv_filename

        except (NotFoundError, ValidationError):
            raise
        except Exception as e:
            logger.error(
                "質問データダウンロード中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                question_type=question_type,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "質問データのダウンロードに失敗しました",
                details={"error": str(e)},
            ) from e

    @measure_performance
    @async_timeout(seconds=300)  # 5分タイムアウト
    async def upload_ppt(self, package: str, phase: str, template: str, file_name: str, file_data: bytes) -> tuple[str, int]:
        """PPTファイルをアップロードします。

        Args:
            package: パッケージ名
            phase: フェーズ名
            template: テンプレート名
            file_name: ファイル名
            file_data: ファイルのバイナリデータ

        Returns:
            tuple[str, int]: (ストレージパス, ファイルサイズ) のタプル

        Raises:
            ValidationError: アップロードに失敗した場合

        Example:
            >>> with open("presentation.pptx", "rb") as f:
            ...     file_data = f.read()
            >>> path, size = await service.upload_ppt(
            ...     "pkg", "phase1", "temp_a", "presentation.pptx", file_data
            ... )
            >>> print(f"Uploaded to {path}, size: {size}")
            Uploaded to pkg/phase1/temp_a/presentation.pptx, size: 1024000

        Note:
            - タイムアウトは5分です
            - 既存のファイルは上書きされます
        """
        logger.info(
            "PPTファイルアップロード中",
            package=package,
            phase=phase,
            template=template,
            file_name=file_name,
            file_size=len(file_data),
            action="upload_ppt",
        )

        try:
            # ストレージパスを生成
            file_path = self._generate_file_path(package, phase, template, file_name)

            # ファイルをアップロード
            await self.storage.upload(self.container, file_path, file_data)

            logger.info(
                "PPTファイルをアップロードしました",
                file_path=file_path,
                file_size=len(file_data),
            )

            return file_path, len(file_data)

        except Exception as e:
            logger.error(
                "PPTファイルアップロード中にエラーが発生しました",
                package=package,
                phase=phase,
                template=template,
                file_name=file_name,
                error=str(e),
                exc_info=True,
            )
            raise ValidationError(
                "PPTファイルのアップロードに失敗しました",
                details={"error": str(e)},
            ) from e
